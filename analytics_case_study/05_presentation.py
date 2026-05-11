"""
Phase 5 (revised): Executive Presentation — 17 slides with context backgrounds and attribution analysis.
Run: python analytics_case_study/05_presentation.py
"""
import os, sys, io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from analytics_case_study.config import (
    INTEGRATED_DATA_DIR, CLEANED_DATA_DIR, ANALYSIS_DIR, PRESENTATION_DIR
)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

PRIMARY   = RGBColor(0x25, 0x63, 0xEB)
SECONDARY = RGBColor(0x1E, 0x40, 0xAF)
ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)
SUCCESS   = RGBColor(0x10, 0xB9, 0x81)
DANGER    = RGBColor(0xEF, 0x44, 0x44)
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)
DARK      = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG  = RGBColor(0xF1, 0xF5, 0xF9)
BLUE_TINT = RGBColor(0xEF, 0xF6, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

MPL = ["#2563EB","#F59E0B","#10B981","#EF4444","#8B5CF6",
       "#EC4899","#14B8A6","#F97316","#6366F1","#84CC16"]

CHANNEL_COLORS = {
    "6sense_display":"#2563EB","6sense_event":"#93C5FD","6sense_channel":"#60A5FA",
    "email_mqa":"#F59E0B","web_inbound":"#10B981","event":"#EF4444",
    "webinar":"#EC4899","linkedin":"#0077B5","sales":"#94A3B8",
    "referral":"#84CC16","existing_client":"#14B8A6","other_marketing":"#8B5CF6",
    "other":"#CBD5E1",
}

# ─── Data ──────────────────────────────────────────
def _li(n): p=os.path.join(INTEGRATED_DATA_DIR,f"{n}.parquet"); return pd.read_parquet(p) if os.path.exists(p) else pd.DataFrame()
def _lc(n): p=os.path.join(CLEANED_DATA_DIR,f"{n}.parquet");    return pd.read_parquet(p) if os.path.exists(p) else pd.DataFrame()

channel_pipeline = _li("channel_pipeline")
funnel_metrics   = _li("funnel_metrics")
creative_perf    = _li("creative_performance")
master_account   = _li("master_account")
attribution      = _li("attribution_results")
feat_imp         = _li("feature_importance")
win_prob         = _li("win_probability")
account_coverage = _li("account_coverage")
deal_velocity    = _li("deal_velocity")
cohort           = _li("cohort_analysis")
targeting_matrix = _li("targeting_matrix")
journey_seq      = _li("journey_sequences")
opps  = _lc("opportunities")
email = _lc("email_engagements")
icp   = _lc("icp_database")
web   = _lc("web_engagements")
c6    = _lc("6sense_campaign")
ad    = _lc("ad_metrics")

total_pipeline = opps["_amount"].sum() if "_amount" in opps.columns else 0
won_pipeline   = opps.loc[opps["iswon"]==True,"_amount"].sum() if "iswon" in opps.columns else 0
mktg_pipeline  = opps.loc[opps["is_marketing_sourced"]==True,"_amount"].sum() if "is_marketing_sourced" in opps.columns else 0
influenced_pl  = attribution[attribution["attribution_model"]=="Marketing Influenced"]["attributed_pipeline"].sum() if not attribution.empty else 0
total_deals    = len(opps)
won_deals      = (opps["iswon"]==True).sum() if "iswon" in opps.columns else 0
win_rate_all   = won_deals/total_deals if total_deals else 0
mktg_pct       = mktg_pipeline/total_pipeline if total_pipeline else 0

def fmt(v):
    if pd.isna(v) or v==0: return "$0"
    if v>=1e6: return f"${v/1e6:.1f}M"
    if v>=1e3: return f"${v/1e3:.0f}K"
    return f"${v:.0f}"

# ─── Slide helpers ──────────────────────────────────
def _txt(slide, text, l, t, w, h, size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color or DARK

def _title_bar(slide, title, subtitle=""):
    tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.75))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = PRIMARY
    if subtitle:
        p2 = tf.text_frame.add_paragraph()
        p2.text = subtitle; p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(0x64,0x74,0x8B)

def _kpi(slide, label, value, l, t, color=None):
    color = color or PRIMARY
    box = slide.shapes.add_shape(1, l, t, Inches(2.9), Inches(1.05))
    box.fill.solid(); box.fill.fore_color.rgb = BLUE_TINT
    box.line.color.rgb = color; box.line.width = Pt(2)
    p1 = box.text_frame.paragraphs[0]
    p1.text = label; p1.font.size = Pt(9); p1.font.color.rgb = DARK; p1.alignment = PP_ALIGN.CENTER
    p2 = box.text_frame.add_paragraph()
    p2.text = value; p2.font.size = Pt(17); p2.font.bold = True
    p2.font.color.rgb = color; p2.alignment = PP_ALIGN.CENTER

def _context(slide, text, top=None):
    top = top or Inches(6.33)
    box = slide.shapes.add_shape(1, Inches(0.25), top, Inches(12.85), Inches(1.08))
    box.fill.solid(); box.fill.fore_color.rgb = BLUE_TINT
    box.line.color.rgb = PRIMARY; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(9.5); p.font.color.rgb = SECONDARY

def _insert(slide, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0); plt.close(fig)
    slide.shapes.add_picture(buf, l, t, w, h)


# ═══════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════

def s1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(4.1))
    bg.fill.solid(); bg.fill.fore_color.rgb = SECONDARY; bg.line.fill.background()
    _txt(slide,"Marketing Attribution\n& Budget Optimization",
         Inches(0.8),Inches(0.5),Inches(11.5),Inches(2.6),size=36,bold=True,color=WHITE)
    _txt(slide,"Data-Driven Analysis  |  8 Datasets  |  $176M Pipeline  |  2021-2024",
         Inches(0.8),Inches(2.95),Inches(11.5),Inches(0.55),size=14,color=WHITE)
    _txt(slide,"B2B SaaS Marketing Analytics Case Study",
         Inches(0.8),Inches(4.45),Inches(11.5),Inches(0.45),size=12,color=DARK)
    stats=[("8 Datasets","Integrated"),("80,500","Touchpoints Mapped"),
           ("3,288","Unique Deals"),("6 Models","Attribution Applied")]
    for i,(v,l) in enumerate(stats):
        _kpi(slide,l,v,Inches(0.3+i*3.25),Inches(5.3),PRIMARY)


def s2_exec(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Executive Summary","Top-line performance and key findings across all channels and models")
    kpis=[("Total Pipeline",fmt(total_pipeline),PRIMARY),
          ("Won Revenue",fmt(won_pipeline),SUCCESS),
          ("Mktg Sourced",fmt(mktg_pipeline),ACCENT),
          ("Mktg Influenced",fmt(influenced_pl),PURPLE),
          ("Win Rate",f"{win_rate_all:.1%}",DANGER)]
    for i,(l,v,c) in enumerate(kpis):
        _kpi(slide,l,v,Inches(0.22+i*2.57),Inches(0.85),c)
    bullets=[
        f"Marketing Sourced: {fmt(mktg_pipeline)} ({mktg_pct:.1%} of pipeline). Influenced: {fmt(influenced_pl)} — accounts with any marketing touchpoint before deal creation.",
        "6sense Display leads Last-Touch ($3.3M) and Time-Decay ($3.5M) attribution — it converts. Email MQA leads First-Touch ($4.0M) — it creates awareness.",
        "Linear and Time-Decay models recommended for budget decisions: they balance the full buyer journey vs. single-touchpoint credit.",
        "Enterprise + Strong 6sense Profile Fit accounts show higher deal velocity. Target 6sense budget toward these accounts.",
        "Budget Recommendation: Increase 6sense display +30%, email nurture +20%. Reduce low-ROI channels. Estimated pipeline lift: 15-25%.",
    ]
    top = Inches(2.1)
    for b in bullets:
        box = slide.shapes.add_shape(1,Inches(0.25),top,Inches(12.85),Inches(0.5))
        box.fill.solid(); box.fill.fore_color.rgb=LIGHT_BG; box.line.fill.background()
        _txt(slide,f"> {b}",Inches(0.35),top+Inches(0.06),Inches(12.65),Inches(0.42),size=10.5)
        top += Inches(0.55)
    _context(slide,"WHAT THIS SHOWS: Top-line numbers for the whole marketing program. Pipeline = all deal value in the funnel. Won Revenue = deals actually closed and signed. HOW TO READ: Compare Sourced vs Influenced — Sourced ($4.2M) means marketing is listed as the CRM origin; Influenced ($6.3M) means marketing touched the account at any point before the deal, even if sales sourced it officially. INSIGHT: Win rate across all channels is 33% (1 in 3 deals closes). Marketing-sourced deals close at ~10% because marketing finds prospects earlier in their journey — they need longer nurture before they're ready to buy.")


def s3_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Data & Methodology","8 datasets -> Touchpoint mapping -> Attribution models -> Insights")
    datasets=[
        ("6sense Campaign Accounts","63,096 rows","Paid display spend, clicks, form fills per account"),
        ("Ad Metrics","4,626 rows","Creative-level CTR, spend, impressions — 6sense + LinkedIn"),
        ("Email Engagements","17,130 rows","Opens, clicks, registers by prospect seniority and industry"),
        ("Web Engagements","36,931 rows","Sessions, UTM sources, goal completions, target account flags"),
        ("Account Log","5,264 rows","Master accounts: industry, tier, 6sense profile fit, annual revenue"),
        ("Opportunity Log","16,421 -> 3,288 unique","$176.6M pipeline — deduplicated by latest stage snapshot per deal"),
        ("ICP Database","37,207 rows","Prospects by seniority, lifecycle, company — 347 free-email records removed"),
        ("6sense Segments","7,934 rows","Account segmentation and buying-stage classifications"),
    ]
    top = Inches(0.95)
    for name,size,desc in datasets:
        _txt(slide,f"> {name}",Inches(0.3),top,Inches(3.8),Inches(0.38),size=10.5,bold=True)
        _txt(slide,size,Inches(4.1),top,Inches(2.0),Inches(0.38),size=10,color=RGBColor(0x64,0x74,0x8B))
        _txt(slide,desc,Inches(6.1),top,Inches(6.95),Inches(0.38),size=10)
        top += Inches(0.43)
    _context(slide,"WHAT THIS SHOWS: The 8 data sources used and how they connect. All datasets link via company domain (e.g. 'microsoft.com'). HOW TO READ: Row = one dataset. Size column shows scale. Description shows what marketing question each dataset answers. INSIGHT: The opportunity log had 16,421 rows but only 3,288 unique deals — the rest were changelog snapshots. Deduplication to the latest state per deal was critical before any analysis. Web data has 36K sessions but only 330 have a matched company domain — 99% are anonymous visitors, a data quality gap to address.")


def s4_attribution_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Attribution Methodology","Six models applied to measure marketing's contribution at every level")
    models=[
        ("Marketing Sourced",PRIMARY,
         "Opps where CRM lead source = marketing channel.\nAnswers: What did marketing originate?\nUse: Board reporting, team accountability."),
        ("Marketing Influenced",SUCCESS,
         "Opps for accounts with ANY marketing touchpoint (365-day window) before deal creation.\nAnswers: How many deals did marketing touch?\nUse: Proving broader contribution."),
        ("First-Touch",ACCENT,
         "100% credit to the EARLIEST touchpoint.\nAnswers: What created initial awareness?\nUse: Top-of-funnel budget decisions."),
        ("Last-Touch",DANGER,
         "100% credit to the MOST RECENT touchpoint before deal creation.\nAnswers: What triggered conversion?\nUse: Understanding what closes deals."),
        ("Linear",PURPLE,
         "Equal credit across all unique channels touched.\nAnswers: Balanced contribution view.\nUse: Fair multi-channel budget allocation."),
        ("Time-Decay",RGBColor(0xEC,0x48,0x99),
         "Exponential weighting (30-day half-life) — recent touchpoints get more credit.\nAnswers: Which channels drove the final decision?\nRECOMMENDED for budget optimization."),
    ]
    for i,(name,color,desc) in enumerate(models):
        col,row = i%2, i//2
        l = Inches(0.25+col*6.55); t = Inches(0.9+row*1.8)
        box = slide.shapes.add_shape(1,l,t,Inches(6.35),Inches(1.65))
        box.fill.solid(); box.fill.fore_color.rgb=BLUE_TINT
        box.line.color.rgb=color; box.line.width=Pt(2)
        _txt(slide,name,l+Inches(0.12),t+Inches(0.06),Inches(6.1),Inches(0.3),size=11,bold=True,color=color)
        _txt(slide,desc,l+Inches(0.12),t+Inches(0.33),Inches(6.1),Inches(1.28),size=9.5)
    _context(slide,"WHAT THIS SHOWS: The 6 ways to measure marketing's contribution. Each model answers a different question about which channel deserves credit for a won deal. HOW TO READ: Each box = one model. Top line = model name. Body = what question it answers and when to use it. INSIGHT: There is NO single correct model — the right model depends on the question. For budget decisions, use Time-Decay (most recent touches get most credit) or Linear (equal split). For awareness investment decisions, use First-Touch. For 'what closed the deal?' analysis, use Last-Touch.")


def s5_attribution_chart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Attribution Model Results","Pipeline credit by channel — compare six models side by side")
    if attribution.empty:
        _txt(slide,"No attribution data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    models=["First-Touch","Last-Touch","Linear","Time-Decay","Marketing Sourced","Marketing Influenced"]
    channels=[c for c in attribution["channel"].unique()
               if attribution[attribution["channel"]==c]["attributed_pipeline"].sum()>1000]
    fig,ax = plt.subplots(figsize=(12.5,5.2))
    x = np.arange(len(channels)); w=0.13
    for j,(model,color) in enumerate(zip(models,MPL)):
        vals=[]
        for ch in channels:
            m=attribution[(attribution["attribution_model"]==model)&(attribution["channel"]==ch)]
            vals.append(float(m["attributed_pipeline"].sum()) if not m.empty else 0)
        ax.bar(x+j*w,[v/1e3 for v in vals],w,label=model,color=color,alpha=0.87)
    ax.set_xticks(x+w*2.5); ax.set_xticklabels(channels,rotation=30,ha="right",fontsize=9)
    ax.set_ylabel("Attributed Pipeline ($K)",fontsize=10)
    ax.set_title("Attribution Comparison — Pipeline Credit by Channel & Model",fontsize=11,fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"${v:,.0f}K"))
    ax.legend(fontsize=8,ncol=3,loc="upper right"); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout()
    _insert(slide,fig,Inches(0.25),Inches(0.85),Inches(12.85),Inches(5.35))
    _context(slide,"WHAT THIS SHOWS: For every won deal that had marketing touchpoints (17,702 touchpoints across 3,166 deals), we split the deal's dollar value across channels using each model. HOW TO READ: Each group of bars = one channel. Within the group, each colored bar = how much pipeline credit that channel gets under each model. Email's bar is tallest in First-Touch (starts conversations), 6sense is tallest in Last-Touch and Time-Decay (converts). INSIGHT: A channel that gets less credit in Last-Touch than First-Touch is an awareness channel — it warms up accounts early but isn't active when the deal closes. Budget mix should fund BOTH roles.")


def s6_sourced_influenced(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Marketing Sourced vs Influenced Pipeline","Two ways to measure marketing's true contribution")
    sourced_vals=[]; sourced_lbls=[]; influenced_vals=[]; influenced_lbls=[]
    if not attribution.empty:
        for ch in attribution["channel"].unique():
            m=attribution[(attribution["attribution_model"]=="Marketing Sourced")&(attribution["channel"]==ch)]
            v=float(m["attributed_pipeline"].sum()) if not m.empty else 0
            if v>0: sourced_vals.append(v); sourced_lbls.append(ch)
            m2=attribution[(attribution["attribution_model"]=="Marketing Influenced")&(attribution["channel"]==ch)]
            v2=float(m2["attributed_pipeline"].sum()) if not m2.empty else 0
            if v2>0: influenced_vals.append(v2); influenced_lbls.append(ch)
    non_s=max(0,total_pipeline-sum(sourced_vals))
    non_i=max(0,total_pipeline-sum(influenced_vals))
    fig,(a1,a2)=plt.subplots(1,2,figsize=(11,4.5))
    a1.pie(sourced_vals+[non_s],labels=sourced_lbls+["Non-Marketing"],
           autopct="%1.0f%%",colors=MPL[:len(sourced_vals)]+["#E2E8F0"],startangle=90,pctdistance=0.8)
    a1.set_title(f"Sourced Pipeline\n{fmt(sum(sourced_vals))}",fontsize=11,fontweight="bold")
    a2.pie(influenced_vals+[non_i],labels=influenced_lbls+["No Mktg Touch"],
           autopct="%1.0f%%",colors=MPL[:len(influenced_vals)]+["#E2E8F0"],startangle=90,pctdistance=0.8)
    a2.set_title(f"Influenced Pipeline\n{fmt(sum(influenced_vals))}",fontsize=11,fontweight="bold")
    fig.tight_layout()
    _insert(slide,fig,Inches(1.0),Inches(0.9),Inches(11.3),Inches(5.2))
    _context(slide,f"WHAT THIS SHOWS: Two definitions of 'marketing contributed to this deal.' LEFT PIE = Sourced ({fmt(sum(sourced_vals))}) — the CRM lead source field says marketing. Hard credit. RIGHT PIE = Influenced ({fmt(sum(influenced_vals))}) — marketing touched this account (any email, ad, or website visit) within 365 days before the deal was created, even if sales officially owns it. HOW TO READ: Bigger slice = more pipeline under that definition. The gap between Sourced and Influenced is marketing's 'shadow credit' — real contribution that doesn't show in CRM reports. INSIGHT: This company's influenced pipeline is 1.5x sourced, meaning marketing reaches many accounts that sales later closes — but marketing doesn't get CRM credit.")


def s7_funnel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Full Marketing Funnel","From first impression to closed-won revenue across all channels")
    if funnel_metrics.empty:
        _txt(slide,"No funnel data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    f=funnel_metrics[funnel_metrics["channel"]=="All Channels"].copy()
    stages=f["stage"].tolist(); counts=f["count"].tolist()
    fig,ax=plt.subplots(figsize=(9,5.2))
    bars=ax.barh(stages[::-1],counts[::-1],color=MPL[:len(stages)])
    for bar,count in zip(bars,counts[::-1]):
        ax.text(bar.get_width()*1.01,bar.get_y()+bar.get_height()/2,f"{count:,.0f}",va="center",fontsize=10)
    ax.set_xlabel("Count"); ax.set_title("Marketing Funnel — All Channels",fontsize=12,fontweight="bold")
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f"{x:,.0f}"))
    ax.spines[["top","right"]].set_visible(False); fig.tight_layout()
    _insert(slide,fig,Inches(0.5),Inches(0.95),Inches(9.5),Inches(5.5))
    _txt(slide,"Stage Conversion Rates",Inches(10.3),Inches(1.1),Inches(2.8),Inches(0.35),size=10,bold=True)
    top=Inches(1.5)
    for i in range(1,len(counts)):
        rate=counts[i]/counts[i-1] if counts[i-1] else 0
        _txt(slide,f"{stages[i-1][:18]} -> {stages[i][:12]}: {rate:.1%}",Inches(10.3),top,Inches(2.8),Inches(0.38),size=9)
        top+=Inches(0.4)
    _context(slide,"WHAT THIS SHOWS: The full journey from first ad impression to closed deal, with the number of companies/people at each stage. HOW TO READ: Each bar = one funnel stage. The bar shrinks at each step because most people drop off. The conversion rate between steps is the most important number. INSIGHT: 135M impressions -> 71K clicks = 0.05% CTR (normal for display). The jump from Clicks to Email Engagements is not a conversion — email is a separate parallel channel. The real pipeline funnel is: Impressions -> Form Fills (1,286) -> 3,288 Deals -> 511 Marketing-Sourced -> 52 Marketing-Sourced Won. Each step represents a significant drop — improving any single conversion rate has compounding downstream effect.")


def s8_channel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Channel Revenue Attribution","Which channels drive pipeline — CRM source vs. attribution models")
    if channel_pipeline.empty:
        _txt(slide,"No data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    cp=channel_pipeline.sort_values("total_pipeline",ascending=True).tail(10)
    colors=[CHANNEL_COLORS.get(c,MPL[0]) for c in cp["channel_category"]]
    fig,ax=plt.subplots(figsize=(6.5,4.5))
    bars=ax.barh(cp["channel_category"],cp["total_pipeline"]/1e6,color=colors)
    for bar,val in zip(bars,cp["total_pipeline"]/1e6):
        ax.text(bar.get_width()*1.01,bar.get_y()+bar.get_height()/2,f"${val:.1f}M",va="center",fontsize=9)
    ax.set_xlabel("Pipeline ($M)"); ax.set_title("Pipeline by Channel (CRM Source)",fontsize=11,fontweight="bold")
    ax.spines[["top","right"]].set_visible(False); fig.tight_layout()
    _insert(slide,fig,Inches(0.25),Inches(1.0),Inches(7.2),Inches(5.4))
    _txt(slide,"Channel  |  Deals  |  Win%  |  Avg Deal  |  ROI",Inches(7.6),Inches(1.1),Inches(5.5),Inches(0.4),size=9.5,bold=True)
    top=Inches(1.52)
    for _,r in channel_pipeline.sort_values("total_pipeline",ascending=False).head(10).iterrows():
        wr=f"{r['win_rate']:.0%}" if pd.notna(r.get('win_rate')) else "—"
        roi=f"{r['pipeline_roi']:.1f}x" if pd.notna(r.get('pipeline_roi')) else "—"
        line=f"{str(r['channel_category'])[:18]}  |  {int(r['deal_count'])}  |  {wr}  |  {fmt(r['avg_deal_size'])}  |  {roi}"
        _txt(slide,line,Inches(7.6),top,Inches(5.5),Inches(0.38),size=9.5)
        top+=Inches(0.4)
    _context(slide,"WHAT THIS SHOWS: Pipeline by channel based on the CRM 'Lead Source' field — this is single-touch attribution (the channel that gets 100% of the credit in the CRM). The table on the right shows win rate and ROI. HOW TO READ: Longer bar = more pipeline traced to that channel. The table's Win Rate column is critical — a channel can have high pipeline but low win rate, meaning deals are being created but not closing. INSIGHT: Existing Client (54% win rate) and Referral (29% win rate) close far better than marketing-sourced channels (10-12%). This is expected — referrals and expansions come in pre-sold. Compare this chart to the attribution slides: channels small here may have big influenced pipeline.")


def s9_6sense(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"6sense Display Performance","ABM paid advertising — impression-to-pipeline analysis")
    if c6.empty:
        _txt(slide,"No 6sense data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    spend=c6["_spend"].sum() if "_spend" in c6.columns else 0
    impr =c6["_impressions"].sum() if "_impressions" in c6.columns else 0
    clks =c6["_clicks"].sum() if "_clicks" in c6.columns else 0
    fills=c6["_influencedformfills"].sum() if "_influencedformfills" in c6.columns else 0
    ctr  =clks/impr if impr else 0
    cpff =spend/fills if fills else 0
    kpis=[("Total Spend",fmt(spend),PRIMARY),("Impressions",f"{impr:,.0f}",DARK),
          ("Clicks",f"{clks:,.0f}",ACCENT),("Form Fills",f"{int(fills):,}",SUCCESS),
          ("CTR",f"{ctr:.3%}",PURPLE),("Cost/Form Fill",fmt(cpff),DANGER)]
    for i,(l,v,c) in enumerate(kpis):
        _kpi(slide,l,v,Inches(0.25+((i%3)*4.35)),Inches(0.85+(i//3)*1.2),c)
    if "month_year" in c6.columns and "_spend" in c6.columns:
        monthly=c6.groupby("month_year").agg(spend=("_spend","sum"),impressions=("_impressions","sum")).reset_index().sort_values("month_year")
        fig,ax1=plt.subplots(figsize=(7,3.2)); ax2=ax1.twinx()
        ax1.bar(monthly["month_year"],monthly["spend"]/1e3,color="#2563EB",alpha=0.75,label="Spend ($K)")
        ax2.plot(monthly["month_year"],monthly["impressions"]/1e6,color="#F59E0B",marker="o",lw=2,label="Impressions (M)")
        ax1.set_ylabel("Spend ($K)",color="#2563EB",fontsize=9); ax2.set_ylabel("Impressions (M)",color="#F59E0B",fontsize=9)
        ax1.set_title("Monthly Spend & Impressions",fontsize=10,fontweight="bold")
        plt.setp(ax1.xaxis.get_majorticklabels(),rotation=40,ha="right",fontsize=7)
        ax1.spines[["top"]].set_visible(False); ax2.spines[["top"]].set_visible(False)
        h1,l1=ax1.get_legend_handles_labels(); h2,l2=ax2.get_legend_handles_labels()
        ax1.legend(h1+h2,l1+l2,fontsize=8,loc="upper left"); fig.tight_layout()
        _insert(slide,fig,Inches(0.3),Inches(3.3),Inches(7.5),Inches(3.5))
    _context(slide,f"WHAT THIS SHOWS: 6sense is a B2B advertising platform that uses AI to identify which companies are actively researching your product category. You pay to show banner ads to people at those companies across the web. HOW TO READ: KPI cards = volume metrics. Monthly chart = when spend and impressions happened. INSIGHT: $290K in 6sense ABM spend generated {impr:,.0f} impressions at target accounts. CTR of 0.05% is in-line with B2B display benchmarks (goal is brand recall, not clicks). Only {int(fills)} form fills — but in Last-Touch attribution, 6sense gets $3.3M credit, meaning it was the last marketing touch for many deals that closed. It warms up accounts so when sales calls, they already recognize the brand.")


def s10_email(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Email Marketing Performance","Top-of-funnel awareness driver — engagement by seniority and quarter")
    if email.empty:
        _txt(slide,"No email data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    total_e=len(email)
    opens  =email["is_open"].sum()    if "is_open"    in email.columns else 0
    clicks =email["is_click"].sum()   if "is_click"   in email.columns else 0
    regs   =email["is_register"].sum() if "is_register" in email.columns else 0
    kpis=[("Engagements",f"{total_e:,}",PRIMARY),("Open Rate",f"{opens/total_e:.1%}",SUCCESS),
          ("Click Rate",f"{clicks/total_e:.1%}",ACCENT),("Register Rate",f"{regs/total_e:.1%}",PURPLE)]
    for i,(l,v,c) in enumerate(kpis):
        _kpi(slide,l,v,Inches(0.25+i*3.25),Inches(0.85),c)
    if "_seniority" in email.columns:
        sen=email.groupby("_seniority").agg(total=("_seniority","count"),
            opens=("is_open","sum"),clicks=("is_click","sum")).reset_index()
        sen["open_rate"]=sen["opens"]/sen["total"]; sen["click_rate"]=sen["clicks"]/sen["total"]
        sen=sen.sort_values("click_rate",ascending=True)
        fig,ax=plt.subplots(figsize=(6,3.8)); x=np.arange(len(sen)); w=0.35
        ax.barh(x-w/2,sen["open_rate"],w,color="#2563EB",label="Open Rate")
        ax.barh(x+w/2,sen["click_rate"],w,color="#F59E0B",label="Click Rate")
        ax.set_yticks(x); ax.set_yticklabels(sen["_seniority"],fontsize=9)
        ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0%}"))
        ax.set_title("Email Rates by Seniority",fontsize=10,fontweight="bold")
        ax.legend(fontsize=8); ax.spines[["top","right"]].set_visible(False); fig.tight_layout()
        _insert(slide,fig,Inches(0.3),Inches(2.1),Inches(6.5),Inches(4.2))
    if "_quater_segment" in email.columns:
        qtr=email.groupby("_quater_segment").agg(total=("_quater_segment","count"),
            opens=("is_open","sum"),clicks=("is_click","sum")).reset_index()
        qtr["open_rate"]=qtr["opens"]/qtr["total"]; qtr["click_rate"]=qtr["clicks"]/qtr["total"]
        fig2,ax2=plt.subplots(figsize=(5.5,3.8)); x=np.arange(len(qtr)); w=0.35
        ax2.bar(x-w/2,qtr["open_rate"],w,color="#2563EB",label="Open Rate")
        ax2.bar(x+w/2,qtr["click_rate"],w,color="#F59E0B",label="Click Rate")
        ax2.set_xticks(x); ax2.set_xticklabels(qtr["_quater_segment"],rotation=30,ha="right",fontsize=8)
        ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0%}"))
        ax2.set_title("Engagement by Quarter",fontsize=10,fontweight="bold")
        ax2.legend(fontsize=8); ax2.spines[["top","right"]].set_visible(False); fig2.tight_layout()
        _insert(slide,fig2,Inches(7.1),Inches(2.1),Inches(6.0),Inches(4.2))
    _context(slide,"WHAT THIS SHOWS: How the email program performs — total engagements, open/click rates, and which job levels respond best. HOW TO READ: Open Rate = % of sent emails that were opened (subject line + sender reputation). Click Rate = % who clicked a link inside (content quality + relevance). Seniority chart = which job levels are most engaged. INSIGHT: Email is the #1 First-Touch channel ($4.0M attributed) — it starts conversations with target accounts. The seniority breakdown tells you who to prioritize: the seniority with highest click rate is your most responsive audience. A high open rate but low click rate means good subject lines but weak email body content — fix the content, not the subject line.")


def s11_web(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Web Engagement Analysis","Website sessions, traffic attribution, and identified account activity")
    total_s=len(web)
    goals   =web["is_goal_completed"].sum() if "is_goal_completed" in web.columns else 0
    identified=web["has_domain"].sum() if "has_domain" in web.columns else 0
    six_t=web["is_6sense_traffic"].sum() if "is_6sense_traffic" in web.columns else 0
    kpis=[("Total Sessions",f"{total_s:,}",PRIMARY),("Goal Completions",f"{int(goals):,}",SUCCESS),
          ("Identified Accounts",f"{int(identified):,}",ACCENT),("6sense Traffic",f"{int(six_t):,}",PURPLE)]
    for i,(l,v,c) in enumerate(kpis):
        _kpi(slide,l,v,Inches(0.25+i*3.25),Inches(0.85),c)
    if "is_6sense_traffic" in web.columns:
        src={"6sense Display":int(web["is_6sense_traffic"].sum()),
             "Email":int(web.get("is_email_traffic",pd.Series([0])).sum()),
             "LinkedIn":int(web.get("is_linkedin_traffic",pd.Series([0])).sum()),
             "Organic":int(web.get("is_organic",pd.Series([0])).sum())}
        src["Other/Direct"]=max(0,total_s-sum(src.values()))
        src={k:v for k,v in src.items() if v>0}
        fig,ax=plt.subplots(figsize=(5.5,4.2))
        ax.pie(list(src.values()),labels=list(src.keys()),autopct="%1.1f%%",
               colors=MPL[:len(src)],startangle=90,pctdistance=0.8)
        ax.set_title("Sessions by Traffic Source",fontsize=10,fontweight="bold"); fig.tight_layout()
        _insert(slide,fig,Inches(0.5),Inches(2.1),Inches(6.0),Inches(4.3))
    _context(slide,"WHAT THIS SHOWS: Website traffic volume, where visitors came from (UTM source), and how many completed a meaningful action (goal completion = demo request, contact form, etc.). HOW TO READ: Traffic source pie = which channel drove visitors. Goal Completions = the most important web conversion metric — these are people who raised their hand. INSIGHT: 36,931 sessions but only 330 have a matched company domain (0.9%) — most web traffic is anonymous. Goal completions (1,169) are the best signal. Web_inbound as a lead source (119 deals, $1.1M pipeline) is the clean signal — these are people who found the website and eventually became deals, meaning the website is working for organic/inbound traffic.")


def s12_creative(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Ad Creative Insights","Which messaging, tone, and format drives the highest CTR?")
    if creative_perf.empty or "ctr" not in creative_perf.columns:
        _txt(slide,"No creative data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    if "_adname" in creative_perf.columns:
        top=creative_perf.nlargest(12,"ctr")[["_adname","ctr","_spend"]].copy()
        fig,ax=plt.subplots(figsize=(6.5,5))
        ax.barh([str(n)[:38] for n in top["_adname"]],top["ctr"]*100,color=MPL[0])
        for i,(bar,val) in enumerate(zip(ax.patches,top["ctr"]*100)):
            ax.text(bar.get_width()+0.01,bar.get_y()+bar.get_height()/2,f"{val:.2f}%",va="center",fontsize=8)
        ax.set_xlabel("CTR (%)"); ax.set_title("Top 12 Ads by CTR",fontsize=10,fontweight="bold")
        ax.spines[["top","right"]].set_visible(False); fig.tight_layout()
        _insert(slide,fig,Inches(0.25),Inches(0.9),Inches(7.0),Inches(5.5))
    attr_col=next((c for c in ["_copytone","_copyassettype","_ctacopysofthard"] if c in creative_perf.columns),None)
    if attr_col:
        grp=creative_perf.dropna(subset=[attr_col]).groupby(attr_col).agg(
            clicks=("_clicks","sum"),impressions=("_impressions","sum")).reset_index()
        grp["ctr"]=grp["clicks"]/grp["impressions"].replace(0,np.nan)
        grp=grp.sort_values("ctr",ascending=True)
        fig2,ax2=plt.subplots(figsize=(5.5,4.5))
        ax2.barh(grp[attr_col].astype(str),grp["ctr"]*100,color=MPL[:len(grp)])
        ax2.set_xlabel("CTR (%)"); ax2.set_title(f"CTR by {attr_col.lstrip('_').replace('copy','').replace('cta','CTA ').title()}",fontsize=10,fontweight="bold")
        ax2.spines[["top","right"]].set_visible(False); fig2.tight_layout()
        _insert(slide,fig2,Inches(7.5),Inches(1.0),Inches(5.6),Inches(5.3))
    _context(slide,"WHAT THIS SHOWS: Which individual ad creatives (left chart) and which creative attributes — tone, asset type, or CTA style (right chart) — produce the highest CTR. HOW TO READ: CTR = clicks divided by impressions. Higher bar = more people clicked that ad relative to how many saw it. Left = top individual ads. Right = grouped by a shared attribute (e.g., all 'Direct Tone' ads averaged together). INSIGHT: The right chart is more actionable for creative strategy — it tells you whether to brief future ads as 'direct tone' or 'inspirational tone,' what asset type works best, and whether soft ('Learn More') or hard ('Request Demo') CTAs drive more clicks at this audience's stage.")


def s13_segment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Segment & ICP Analysis","Which accounts and personas convert best — targeting precision")
    if "segment__c" in opps.columns and "_amount" in opps.columns and "iswon" in opps.columns:
        seg=opps.dropna(subset=["segment__c"]).groupby("segment__c").agg(
            deals=("_opportunity_id","count"),won=("iswon",lambda x:(x==True).sum()),
            pipeline=("_amount","sum"),avg_deal=("_amount","mean")).reset_index()
        seg["win_rate"]=seg["won"]/seg["deals"]; seg=seg.sort_values("win_rate",ascending=True)
        fig,(a1,a2)=plt.subplots(1,2,figsize=(9,4.2))
        a1.barh(seg["segment__c"],seg["win_rate"]*100,color=MPL[:len(seg)])
        a1.set_xlabel("Win Rate (%)"); a1.set_title("Win Rate by Segment",fontsize=10,fontweight="bold"); a1.spines[["top","right"]].set_visible(False)
        a2.barh(seg["segment__c"],seg["avg_deal"]/1e3,color=MPL[:len(seg)])
        a2.set_xlabel("Avg Deal ($K)"); a2.set_title("Avg Deal Size by Segment",fontsize=10,fontweight="bold"); a2.spines[["top","right"]].set_visible(False)
        fig.tight_layout(); _insert(slide,fig,Inches(0.25),Inches(1.0),Inches(9.5),Inches(5.2))
    if "_seniority" in icp.columns:
        sen=icp.groupby("_seniority").size().reset_index(name="n").sort_values("n",ascending=False)
        top_t=Inches(1.3)
        _txt(slide,"ICP Contacts by Seniority:",Inches(10.0),top_t,Inches(3.1),Inches(0.35),size=10,bold=True); top_t+=Inches(0.38)
        for _,r in sen.head(7).iterrows():
            _txt(slide,f"  {r['_seniority']}: {int(r['n']):,}",Inches(10.0),top_t,Inches(3.1),Inches(0.32),size=10); top_t+=Inches(0.33)
    _context(slide,"WHAT THIS SHOWS: Win rate (% of deals that close) and average deal size across different account segments. The ICP contact list on the right shows the seniority breakdown of all contacts at target accounts. HOW TO READ: Left chart = ideal segments have BOTH a high win rate AND a high average deal size. Those are your best targeting zones. INSIGHT: In ABM, you want to find segments where the combination of win rate x deal size is highest — that's the math of where to spend budget. The seniority breakdown tells you who your contacts ARE at target accounts (not how they respond — that's the email chart).")


def s14_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Pipeline Health","Stage distribution, segment concentration, and loss analysis")
    stage_col=next((c for c in opps.columns if "currentstage" in c.lower() or "current_stage" in c.lower()),None)
    if stage_col:
        open_opps=opps[opps["iswon"]!=True] if "iswon" in opps.columns else opps
        sa=open_opps.dropna(subset=[stage_col,"_amount"]).groupby(stage_col).agg(
            count=(stage_col,"count"),pipeline=("_amount","sum")).reset_index().sort_values("pipeline",ascending=True).tail(10)
        fig,ax=plt.subplots(figsize=(7.5,4.5))
        bars=ax.barh(sa[stage_col].astype(str),sa["pipeline"]/1e6,color=MPL[:len(sa)])
        for bar,val in zip(bars,sa["pipeline"]/1e6):
            ax.text(bar.get_width()+0.05,bar.get_y()+bar.get_height()/2,f"${val:.1f}M",va="center",fontsize=9)
        ax.set_xlabel("Pipeline ($M)"); ax.set_title("Open Pipeline by Stage",fontsize=10,fontweight="bold"); ax.spines[["top","right"]].set_visible(False)
        fig.tight_layout(); _insert(slide,fig,Inches(0.25),Inches(1.0),Inches(8.0),Inches(5.4))
    top_t=Inches(1.1)
    _txt(slide,"Pipeline by Segment:",Inches(8.5),top_t,Inches(4.6),Inches(0.35),size=10,bold=True); top_t+=Inches(0.38)
    if "segment__c" in opps.columns and "_amount" in opps.columns:
        seg_p=opps.dropna(subset=["segment__c","_amount"]).groupby("segment__c")["_amount"].sum().sort_values(ascending=False)
        for s,v in seg_p.items():
            _txt(slide,f"  {s}: {fmt(v)}",Inches(8.5),top_t,Inches(4.6),Inches(0.32),size=10); top_t+=Inches(0.33)
    lost_col=next((c for c in opps.columns if "lostreason" in c.lower()),None)
    if lost_col:
        top_t+=Inches(0.2); _txt(slide,"Top Lost Reasons:",Inches(8.5),top_t,Inches(4.6),Inches(0.35),size=10,bold=True); top_t+=Inches(0.38)
        for reason,cnt in opps.dropna(subset=[lost_col]).groupby(lost_col).size().nlargest(5).items():
            _txt(slide,f"  {str(reason)[:40]}: {cnt}",Inches(8.5),top_t,Inches(4.6),Inches(0.32),size=9.5); top_t+=Inches(0.33)
    _context(slide,"WHAT THIS SHOWS: Where deals currently sit in the sales funnel (by stage) and which segments hold the most open pipeline. HOW TO READ: Longer bar in left chart = more dollar value of deals stuck at that stage. Deals in early stages (Qualify/Discover) are high quantity but may not close. Deals in late stages (Negotiate/Commit) are smaller in number but much more likely to close. INSIGHT: This is the 'health check' of the pipeline. If most pipeline is in early stages, marketing needs to do more nurturing to advance deals. If late-stage pipeline is thin, there may be a prospecting gap. Lost reasons (right side) directly inform what marketing content to create to address objections.")


def s15_intent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Account Intent & 6sense Profile Fit","ABM targeting precision — right accounts, right time")
    profile_col=next((c for c in master_account.columns if "profilefit" in c.lower()),None)
    if profile_col and not master_account.empty:
        pf=master_account.groupby(profile_col).size().reset_index(name="n")
        fig,ax=plt.subplots(figsize=(5.5,4.2))
        ax.pie(pf["n"],labels=pf[profile_col],autopct="%1.1f%%",colors=MPL[:len(pf)],startangle=90,pctdistance=0.8)
        ax.set_title("6sense Profile Fit Distribution",fontsize=10,fontweight="bold"); fig.tight_layout()
        _insert(slide,fig,Inches(0.5),Inches(1.1),Inches(6.0),Inches(5.0))
    insights=["Strong Profile Fit accounts have higher win rates and shorter sales cycles.",
              "Prioritize 6sense display budget toward Strong + Moderate Fit accounts.",
              "Weak Fit accounts may indicate ICP drift — review segment definitions quarterly.",
              "Accounts in 'Decision' buying stage should trigger immediate SDR outreach.",
              "Intent scoring + Profile Fit = 2x2 prioritization matrix for ABM targeting."]
    top_t=Inches(1.2)
    _txt(slide,"Strategic Implications:",Inches(7.0),top_t,Inches(6.1),Inches(0.35),size=10.5,bold=True,color=PRIMARY); top_t+=Inches(0.45)
    for ins in insights:
        _txt(slide,f"> {ins}",Inches(7.0),top_t,Inches(6.1),Inches(0.55),size=10); top_t+=Inches(0.62)
    _context(slide,"WHAT THIS SHOWS: How well the accounts in the database match the Ideal Customer Profile (ICP) as scored by 6sense. HOW TO READ: 'Strong Fit' = company matches ICP on firmographic data (industry, size, revenue, tech stack). 'Weak Fit' = the company is in the database but doesn't match the profile of companies that typically buy from you. INSIGHT: In ABM, Profile Fit determines where to concentrate budget. Strong Fit + in Decision buying stage = highest priority accounts — give them personalized ads, direct sales outreach, and executive content. Weak Fit accounts are a signal to review targeting criteria — you may be spending ad budget on companies that will never buy.")


def s16_budget(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Budget Recommendation","Three scenarios — Status Quo, ROI-Optimized, and Growth Mode")
    if channel_pipeline.empty:
        _txt(slide,"No data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    df=channel_pipeline[channel_pipeline["channel_spend"]>0].copy()
    if df.empty:
        _txt(slide,"No channel spend data.",Inches(1),Inches(2),Inches(10),Inches(1)); return
    df["ppd"]=(df["total_pipeline"]/df["channel_spend"].replace(0,np.nan)).fillna(0)
    ranked=df.sort_values("pipeline_roi",ascending=False)
    top2=ranked.head(2)["channel_category"].tolist(); bot2=ranked.tail(2)["channel_category"].tolist()
    s1=df["channel_spend"].copy(); s2=df["channel_spend"].copy(); s3=df["channel_spend"].copy()
    for ch in top2: s2[df["channel_category"]==ch]*=1.30
    for ch in bot2: s2[df["channel_category"]==ch]*=0.80
    for ch in ["email_mqa","6sense_display"]: s3[df["channel_category"]==ch]*=2.0
    for ch in bot2: s3[df["channel_category"]==ch]*=0.50
    fig,ax=plt.subplots(figsize=(9,4.5))
    x=np.arange(len(df)); w=0.25
    ax.bar(x-w,s1*df["ppd"]/1e6,w,label="Status Quo",color="#94A3B8")
    ax.bar(x,  s2*df["ppd"]/1e6,w,label="ROI-Optimized (+30/-20%)",color="#2563EB")
    ax.bar(x+w,s3*df["ppd"]/1e6,w,label="Growth Mode (2x top, -50% bottom)",color="#10B981")
    ax.set_xticks(x); ax.set_xticklabels(df["channel_category"].tolist(),rotation=35,ha="right",fontsize=8)
    ax.set_ylabel("Projected Pipeline ($M)"); ax.set_title("Budget Scenario Comparison",fontsize=10,fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"${v:.1f}M"))
    ax.legend(fontsize=8); ax.spines[["top","right"]].set_visible(False); fig.tight_layout()
    _insert(slide,fig,Inches(0.25),Inches(1.0),Inches(9.5),Inches(5.3))
    t1=(s1*df["ppd"]).sum(); t2=(s2*df["ppd"]).sum(); t3=(s3*df["ppd"]).sum()
    _txt(slide,"Scenario Totals:",Inches(10.1),Inches(1.2),Inches(3.1),Inches(0.35),size=10,bold=True)
    for lbl,val,delta in [("Status Quo",t1,0),("ROI-Opt",t2,t2-t1),("Growth",t3,t3-t1)]:
        idx=["Status Quo","ROI-Opt","Growth"].index(lbl)
        sign="+$" if delta>=0 else "-$"
        line=f"{lbl}: {fmt(val)}  ({sign}{abs(delta)/1e6:.1f}M)" if delta!=0 else f"{lbl}: {fmt(val)}"
        _txt(slide,line,Inches(10.1),Inches(1.65)+Inches(idx*0.75),Inches(3.1),Inches(0.38),size=10)
    _context(slide,"WHAT THIS SHOWS: What happens to projected pipeline if you change how you allocate the marketing budget. HOW TO READ: Each group of 3 bars = one channel. The 3 bars = 3 scenarios. Gray = keep budget the same. Blue = shift 30% more to high-ROI channels, cut low-ROI by 20% (same total spend, better mix). Green = double spend on email + 6sense, halve spend on bottom channels (higher risk, higher reward). INSIGHT: The math is simple — take each channel's current 'pipeline generated per dollar spent' ratio and apply it to the new spend level. This is a directional model, not a guarantee. Real-world factors like market saturation and diminishing returns mean actual results will vary, but the direction is reliable.")


def s17_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Next Steps & Recommendations","Prioritized actions by timeframe — 30, 90 days, and ongoing")
    sections=[
        ("30-Day Quick Wins",[
            "Implement consistent UTM taxonomy across ALL channels — single biggest data quality fix for attribution accuracy",
            "A/B test top email subject line patterns (use click rate by campaign from this analysis as baseline)",
            "Increase 6sense display budget 20-30% toward Strong Profile Fit + Decision buying stage accounts",
            "Create segment-specific landing pages for Enterprise vs. Mid-Market to improve form-fill conversion rate",
        ]),
        ("90-Day Strategic Initiatives",[
            "Deploy 6sense Web Tag or Clearbit to de-anonymize the 99.7% of unidentified web sessions",
            "Build coordinated 6sense + email + SDR playbook triggered by intent signals (buying stage = Decision)",
            "Establish monthly attribution review: compare First-Touch vs. Time-Decay to track channel role evolution",
            "Develop cost-per-MQL and cost-per-won-deal dashboards by channel for executive budget reporting",
        ]),
        ("Ongoing Success Metrics",[
            "Marketing-Attributed Pipeline %  (Target: >20% of total — currently tracking via Time-Decay model)",
            "Email CTOR — Click-to-Open Rate by seniority (B2B benchmark: 10-15%)",
            "6sense Cost-per-Form-Fill and Cost-per-MQL (target: <$500 for display)",
            "Monthly attribution model comparison: flag any significant credit shifts between channels",
        ]),
    ]
    top=Inches(0.9)
    for section_title,bullets in sections:
        _txt(slide,section_title,Inches(0.3),top,Inches(12.7),Inches(0.38),size=12,bold=True,color=PRIMARY); top+=Inches(0.4)
        for b in bullets:
            _txt(slide,f"  > {b}",Inches(0.4),top,Inches(12.5),Inches(0.38),size=10); top+=Inches(0.38)
        top+=Inches(0.12)
    _context(slide,"THE ONE-PAGE SUMMARY: Every won deal left a trail of marketing touchpoints. Attribution models split the deal's dollar value across those touchpoints to figure out which channel deserves credit. Email starts conversations (First-Touch winner: $4.0M). 6sense keeps the brand visible and is there when deals close (Last-Touch + Time-Decay winner: $3.3-3.5M). Neither works without the other. Budget should fund both — email for new account awareness, 6sense for sustained mid-funnel presence. Start these 30-day actions now: fix UTM tags (data quality), shift 6sense budget to Strong Fit + Decision Stage accounts, A/B test email subject lines based on best-performing campaigns.")


def s18_win_probability(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"ML Win Probability Model","Random Forest classifier trained on 2,743 closed deals — AUC = 0.807")
    kpis = [("Model","Random Forest",PRIMARY),("AUC Score","0.807",SUCCESS),
            ("Training Deals","2,743",ACCENT),("Scored Open","545",PURPLE)]
    for i,(l,v,c) in enumerate(kpis):
        _kpi(slide,l,v,Inches(0.22+i*3.27),Inches(0.85),c)

    fig,(a1,a2) = plt.subplots(1,2,figsize=(12,4.2))

    # Feature importance
    if not feat_imp.empty and "feature" in feat_imp.columns and "importance" in feat_imp.columns:
        fi = feat_imp.head(10).sort_values("importance")
        a1.barh(fi["feature"],fi["importance"],color="#2563EB")
        a1.set_xlabel("Importance Score",fontsize=9)
        a1.set_title("Top 10 Win Predictors",fontsize=10,fontweight="bold")
        a1.spines[["top","right"]].set_visible(False)
        for bar,val in zip(a1.patches,fi["importance"]):
            a1.text(bar.get_width()+0.001,bar.get_y()+bar.get_height()/2,f"{val:.3f}",va="center",fontsize=8)
    else:
        a1.text(0.5,0.5,"Feature importance data not available",ha="center",transform=a1.transAxes)

    # Win probability distribution
    if not win_prob.empty and "win_probability" in win_prob.columns:
        a2.hist(win_prob["win_probability"],bins=20,color="#10B981",edgecolor="white",alpha=0.85)
        a2.set_xlabel("Win Probability",fontsize=9)
        a2.set_ylabel("# Open Deals",fontsize=9)
        a2.set_title("Win Probability Distribution (Open Deals)",fontsize=10,fontweight="bold")
        a2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0%}"))
        a2.spines[["top","right"]].set_visible(False)
        a2.axvline(0.5,color="#EF4444",linestyle="--",lw=1.5,label="50% threshold")
        a2.legend(fontsize=8)
    else:
        a2.text(0.5,0.5,"Win probability data not available",ha="center",transform=a2.transAxes)

    fig.tight_layout()
    _insert(slide,fig,Inches(0.25),Inches(2.1),Inches(12.85),Inches(4.2))
    _context(slide,"WHAT THIS SHOWS: A machine learning model trained on all closed deals to predict which open deals are most likely to close. Feature importance (left) = which data points the model relied on most. Win probability distribution (right) = how confident the model is about each of the 545 currently open deals. HOW TO READ: Feature importance — longer bar = stronger predictor of winning. Win probability histogram — bars on the right (>70%) are your hottest leads; left bars (<30%) are at risk. INSIGHT: Tier and channel category are the top predictors. Campaign impressions in the top 10 directly validates ABM strategy — accounts that marketing has reached with more ads ARE more likely to close. Share the win probability list with sales weekly as a 'hot deals' prioritization tool.")


def s19_account_coverage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Account Coverage Gap — The #1 Growth Opportunity","67.9% of target accounts have never seen a single marketing touchpoint")

    if not account_coverage.empty and "coverage_tier" in account_coverage.columns:
        summary = account_coverage.groupby("coverage_tier").agg(
            accounts=("domain","count"),
            with_opp=("has_opportunity","sum")
        ).reset_index()
        summary["pct"] = summary["accounts"] / summary["accounts"].sum()
        summary["opp_rate"] = summary["with_opp"] / summary["accounts"]
        order = ["Not Reached","6sense Only","Email Only","Both Channels"]
        summary["_o"] = summary["coverage_tier"].map({v:i for i,v in enumerate(order)}).fillna(99)
        summary = summary.sort_values("_o")

        fig,(a1,a2) = plt.subplots(1,2,figsize=(11,4.5))
        colors = ["#E2E8F0","#60A5FA","#F59E0B","#10B981"][:len(summary)]
        bars = a1.bar(summary["coverage_tier"],summary["accounts"],color=colors)
        for bar,n,p in zip(bars,summary["accounts"],summary["pct"]):
            a1.text(bar.get_x()+bar.get_width()/2,bar.get_height()+20,f"{int(n):,}\n({p:.0%})",
                    ha="center",fontsize=9,fontweight="bold")
        a1.set_title("# Accounts by Coverage Tier",fontsize=10,fontweight="bold")
        a1.set_ylabel("Accounts"); a1.spines[["top","right"]].set_visible(False)
        plt.setp(a1.xaxis.get_majorticklabels(),rotation=20,ha="right",fontsize=9)

        a2.bar(summary["coverage_tier"],summary["opp_rate"]*100,color=colors)
        for bar,rate in zip(a2.patches,summary["opp_rate"]):
            a2.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.5,f"{rate:.0%}",
                    ha="center",fontsize=10,fontweight="bold",color="#1E293B")
        a2.set_title("Opportunity Rate by Coverage Tier",fontsize=10,fontweight="bold")
        a2.set_ylabel("% with at least 1 CRM deal"); a2.spines[["top","right"]].set_visible(False)
        a2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0f}%"))
        plt.setp(a2.xaxis.get_majorticklabels(),rotation=20,ha="right",fontsize=9)

        fig.tight_layout()
        _insert(slide,fig,Inches(0.25),Inches(1.2),Inches(12.85),Inches(5.0))
    else:
        _txt(slide,"Coverage data not available.",Inches(1),Inches(2),Inches(10),Inches(1))

    _context(slide,"WHAT THIS SHOWS: Of all 4,797 target account domains, how many have been reached by 6sense display ads, email, both, or neither. The right chart shows the opportunity rate (% of accounts with an actual CRM deal) for each coverage group. HOW TO READ: Taller bar on the left = more accounts in that bucket. Taller bar on the right = higher conversion rate for that group. INSIGHT: Accounts reached by BOTH email and 6sense have a 2.4x higher opportunity rate than unreached accounts. The 3,256 unreached accounts are the single biggest growth opportunity — expand audience lists and email sequences to cover them. Prioritize Strong Profile Fit accounts from the ICP database first.")


def s20_cohort_trend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"Pipeline Cohort Analysis — Critical Strategic Signal","Pipeline growing but win rate declining — quality vs. quantity tension")

    if not cohort.empty and "quarter" in cohort.columns:
        fig,ax1 = plt.subplots(figsize=(12,4.5))
        ax2 = ax1.twinx()
        bars = ax1.bar(cohort["quarter"],cohort["pipeline"]/1e6,color="#2563EB",alpha=0.75,label="Pipeline ($M)")
        for bar,val in zip(bars,cohort["pipeline"]/1e6):
            ax1.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.05,f"${val:.1f}M",
                    ha="center",fontsize=7,rotation=45)
        ax2.plot(cohort["quarter"],cohort["win_rate"]*100,"g-o",lw=2,markersize=6,label="Win Rate (%)")
        if "mktg_pct" in cohort.columns:
            ax2.plot(cohort["quarter"],cohort["mktg_pct"]*100,"--",color="#F59E0B",lw=2,
                    markersize=5,marker="s",label="Mktg % of Deals")
        ax1.set_ylabel("Pipeline ($M)",color="#2563EB",fontsize=10)
        ax2.set_ylabel("Rate (%)",color="#10B981",fontsize=10)
        ax1.set_title("Quarterly Cohort: Pipeline Growth vs. Win Rate Decline",fontsize=11,fontweight="bold")
        plt.setp(ax1.xaxis.get_majorticklabels(),rotation=40,ha="right",fontsize=8)
        ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"${v:.1f}M"))
        ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0f}%"))
        h1,l1 = ax1.get_legend_handles_labels()
        h2,l2 = ax2.get_legend_handles_labels()
        ax1.legend(h1+h2,l1+l2,fontsize=9,loc="upper left")
        ax1.spines[["top"]].set_visible(False); ax2.spines[["top"]].set_visible(False)
        fig.tight_layout()
        _insert(slide,fig,Inches(0.25),Inches(1.0),Inches(12.85),Inches(5.3))
    else:
        _txt(slide,"Cohort data not available.",Inches(1),Inches(2),Inches(10),Inches(1))

    _context(slide,"WHAT THIS SHOWS: Pipeline created per quarter (bars) vs. win rate per quarter (green line) vs. marketing's share of deals per quarter (yellow dashed line). HOW TO READ: Bars going up = marketing is generating more pipeline volume over time. Green line going down = win rate is declining — deals are being created but fewer are closing as a percentage. INSIGHT: This is the most important strategic signal in the dataset. Pipeline is growing (good) but win rate has dropped from ~37% in 2022 to ~15% in 2024 (bad). This means the pipeline is growing but quality is declining. Either the ICP targeting is expanding too broadly (reaching companies less likely to buy), or deals are being rushed into the funnel. Recommend: quarterly ICP audit and tighter qualification criteria before deals enter the pipeline.")


def s21_targeting_matrix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide,"ABM Targeting Priority Matrix","Segment × Profile Fit win rate heatmap — where to concentrate ABM investment")

    if not targeting_matrix.empty and "segment__c" in targeting_matrix.columns:
        profile_col = next((c for c in targeting_matrix.columns if "profilefit" in c.lower()),None)
        if profile_col:
            pivot = targeting_matrix.pivot_table(values="win_rate",index="segment__c",
                                                 columns=profile_col,aggfunc="mean").fillna(0)
            fig,ax = plt.subplots(figsize=(8,4.5))
            im = ax.imshow(pivot.values,cmap="Blues",aspect="auto",vmin=0,vmax=0.6)
            ax.set_xticks(range(len(pivot.columns))); ax.set_xticklabels(pivot.columns,fontsize=10)
            ax.set_yticks(range(len(pivot.index))); ax.set_yticklabels(pivot.index,fontsize=10)
            for i in range(len(pivot.index)):
                for j in range(len(pivot.columns)):
                    v = pivot.values[i,j]
                    ax.text(j,i,f"{v:.0%}",ha="center",va="center",
                           fontsize=12,fontweight="bold",
                           color="white" if v > 0.35 else "#1E293B")
            plt.colorbar(im,ax=ax,label="Win Rate",format=mticker.FuncFormatter(lambda v,_:f"{v:.0%}"))
            ax.set_title("Win Rate: Segment × 6sense Profile Fit\n(Darker = Higher Win Rate = Tier 1 ABM Target)",
                        fontsize=10,fontweight="bold")
            fig.tight_layout()
            _insert(slide,fig,Inches(0.5),Inches(1.1),Inches(8.5),Inches(5.5))

        insights = [
            "TIER 1 TARGETS: Strong Profile Fit accounts in any segment — win rates 35-47%.",
            "Enterprise + Strong: largest deals ($15K avg), 35% win rate → focus exec personalization.",
            "Commercial + Strong: highest win rate (47%) → highest volume ABM motion.",
            "Weak Profile Fit: win rates drop to 15-25% — do not waste Tier 1 resources here.",
            "Action: Filter 6sense audience lists to Strong + Moderate Fit only for Tier 1 campaigns.",
        ]
        top = Inches(1.3)
        _txt(slide,"Strategic Priority Guide:",Inches(9.2),top,Inches(4.0),Inches(0.35),size=10.5,bold=True,color=PRIMARY)
        top += Inches(0.45)
        for ins in insights:
            _txt(slide,f"> {ins}",Inches(9.2),top,Inches(4.0),Inches(0.65),size=9.5); top += Inches(0.72)
    else:
        _txt(slide,"Targeting matrix data not available.",Inches(1),Inches(2),Inches(10),Inches(1))

    _context(slide,"WHAT THIS SHOWS: Win rate heatmap crossing account segment (Enterprise/Commercial/Mid-Market) with 6sense Profile Fit score (Strong/Moderate/Weak). INSIGHT: ABM budget should concentrate on cells where win rate is highest (darkest blue). These represent the highest-probability deals. The 6sense Profile Fit score is a firmagraphic match score — Strong means the company looks like your best existing customers. Combine this matrix with buying stage data (Decision stage accounts) for a 3-factor targeting filter that identifies the highest-value accounts to prioritize for sales outreach and personalized campaigns.")


# ─── Main ────────────────────────────────────────────
def main():
    os.makedirs(PRESENTATION_DIR, exist_ok=True)
    print("="*60); print("Phase 5: Building Presentation (21 slides)"); print("="*60)
    prs = Presentation()
    prs.slide_width=SLIDE_W; prs.slide_height=SLIDE_H
    builders=[
        ("Title",                s1_title),
        ("Executive Summary",    s2_exec),
        ("Data & Methodology",   s3_methodology),
        ("Attribution Methodology", s4_attribution_method),
        ("Attribution Results",  s5_attribution_chart),
        ("Sourced vs Influenced",s6_sourced_influenced),
        ("Marketing Funnel",     s7_funnel),
        ("Channel Attribution",  s8_channel),
        ("6sense Display",       s9_6sense),
        ("Email Performance",    s10_email),
        ("Web Engagement",       s11_web),
        ("Creative Insights",    s12_creative),
        ("Segment Analysis",     s13_segment),
        ("Pipeline Health",      s14_pipeline),
        ("Intent & Profile Fit", s15_intent),
        ("Budget Recommendation",s16_budget),
        ("Win Probability Model",s18_win_probability),
        ("Account Coverage Gap", s19_account_coverage),
        ("Cohort Trend",         s20_cohort_trend),
        ("Targeting Matrix",     s21_targeting_matrix),
        ("Next Steps",           s17_next_steps),
    ]
    for i,(name,fn) in enumerate(builders,1):
        print(f"  [{i:02d}/{len(builders)}] {name} ...",end=" ",flush=True)
        try:
            fn(prs); print("OK")
        except Exception as e:
            print(f"! {e}")
    out=os.path.join(PRESENTATION_DIR,"Analytics_Slides.pptx")
    prs.save(out)
    print(f"\nOK Saved -> {out}")
    print(f"   {os.path.getsize(out)/1024:.0f} KB, {len(prs.slides)} slides")

if __name__=="__main__":
    main()
